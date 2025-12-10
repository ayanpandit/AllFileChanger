from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
import io
import os
import tempfile

app = Flask(__name__)
CORS(app)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'service': 'powerpoint-converter'})

@app.route('/convert', methods=['POST'])
def convert_ppt():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        target_format = request.form.get('format', 'pdf')
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp:
            file.save(temp.name)
            temp_path = temp.name
        
        if target_format == 'pdf':
            # Note: This requires LibreOffice for conversion
            from subprocess import run
            output_path = temp_path.replace('.pptx', '.pdf')
            run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(temp_path), temp_path])
            mimetype = 'application/pdf'
            filename = 'converted.pdf'
        elif target_format == 'txt':
            prs = Presentation(temp_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text.append(shape.text)
            
            output = io.BytesIO('\n'.join(text).encode('utf-8'))
            os.unlink(temp_path)
            output.seek(0)
            return send_file(output, mimetype='text/plain', as_attachment=True, download_name='converted.txt')
        else:
            os.unlink(temp_path)
            return jsonify({'error': 'Unsupported format'}), 400
        
        with open(output_path, 'rb') as f:
            output = io.BytesIO(f.read())
        
        os.unlink(temp_path)
        os.unlink(output_path)
        
        output.seek(0)
        return send_file(output, mimetype=mimetype, as_attachment=True, download_name=filename)
        
    except Exception as e:
        return jsonify({'error': 'Failed to convert PowerPoint', 'details': str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5021))
    app.run(host='0.0.0.0', port=port, debug=True)
