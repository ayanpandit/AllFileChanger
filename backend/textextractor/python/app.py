from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io

app = Flask(__name__)
CORS(app)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'service': 'text-extractor'})

@app.route('/extract', methods=['POST'])
def extract_text():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        file_type = file.filename.split('.')[-1].lower()
        
        text = []
        
        if file_type == 'pdf':
            reader = PdfReader(io.BytesIO(file.read()))
            for page in reader.pages:
                text.append(page.extract_text())
        
        elif file_type in ['docx', 'doc']:
            doc = Document(io.BytesIO(file.read()))
            for para in doc.paragraphs:
                text.append(para.text)
        
        elif file_type in ['pptx', 'ppt']:
            prs = Presentation(io.BytesIO(file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text.append(shape.text)
        
        elif file_type == 'txt':
            text.append(file.read().decode('utf-8'))
        
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        extracted_text = '\n\n'.join(text)
        output = io.BytesIO(extracted_text.encode('utf-8'))
        output.seek(0)
        
        return send_file(output, mimetype='text/plain', as_attachment=True, download_name='extracted_text.txt')
        
    except Exception as e:
        return jsonify({'error': 'Failed to extract text', 'details': str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5022))
    app.run(host='0.0.0.0', port=port, debug=True)
