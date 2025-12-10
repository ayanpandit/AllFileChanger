from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
import io
import os
import tempfile
from pdf2image import convert_from_path

app = Flask(__name__)
CORS(app)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy', 'service': 'pdf-to-powerpoint'})

@app.route('/convert', methods=['POST'])
def pdf_to_ppt():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file provided'}), 400
        
        pdf_file = request.files['pdf']
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as pdf_temp:
            pdf_file.save(pdf_temp.name)
            pdf_path = pdf_temp.name
        
        # Convert PDF pages to images
        images = convert_from_path(pdf_path)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        for i, image in enumerate(images):
            # Save image temporarily
            img_path = f'{pdf_path}_{i}.png'
            image.save(img_path, 'PNG')
            
            # Add slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add image to slide
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(img_path, left, top, width=prs.slide_width)
            
            # Clean up temp image
            os.unlink(img_path)
        
        # Save PowerPoint
        ppt_path = pdf_path.replace('.pdf', '.pptx')
        prs.save(ppt_path)
        
        # Read the PowerPoint file
        with open(ppt_path, 'rb') as f:
            output = io.BytesIO(f.read())
        
        # Cleanup
        os.unlink(pdf_path)
        os.unlink(ppt_path)
        
        output.seek(0)
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='converted.pptx'
        )
        
    except Exception as e:
        return jsonify({'error': 'Failed to convert PDF to PowerPoint', 'details': str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5016))
    app.run(host='0.0.0.0', port=port, debug=True)
