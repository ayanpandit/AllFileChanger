from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
import io
import os
import tempfile
import subprocess
import logging
from werkzeug.utils import secure_filename
import shutil
from pathlib import Path

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# Production configurations
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
LIBREOFFICE_TIMEOUT = 120  # 2 minutes timeout for conversion

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def check_libreoffice():
    """Check if LibreOffice is available"""
    try:
        result = subprocess.run(
            ['libreoffice', '--version'],
            capture_output=True,
            timeout=5,
            text=True
        )
        return result.returncode == 0
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return False

# Check LibreOffice availability at startup
LIBREOFFICE_AVAILABLE = check_libreoffice()
if not LIBREOFFICE_AVAILABLE:
    logger.warning("LibreOffice not found. PDF conversion will not work.")

@app.route('/health', methods=['GET'])
def health():
    return jsonify({
        'status': 'healthy',
        'service': 'pptx-to-pdf-converter',
        'libreoffice_available': LIBREOFFICE_AVAILABLE
    })

@app.route('/convert', methods=['POST'])
def convert_ppt():
    temp_input = None
    temp_output = None
    temp_dir = None
    
    try:
        # Validate file presence
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'Invalid file type. Only PPT and PPTX files are allowed'}), 400
        
        target_format = request.form.get('format', 'pdf').lower()
        
        # Create a temporary directory for this conversion
        temp_dir = tempfile.mkdtemp()
        
        # Secure the filename and save
        original_filename = secure_filename(file.filename)
        file_extension = original_filename.rsplit('.', 1)[1].lower()
        temp_input = os.path.join(temp_dir, f'input.{file_extension}')
        
        file.save(temp_input)
        logger.info(f"Processing file: {original_filename}, size: {os.path.getsize(temp_input)} bytes")
        
        if target_format == 'pdf':
            if not LIBREOFFICE_AVAILABLE:
                return jsonify({'error': 'PDF conversion is not available. LibreOffice is not installed.'}), 503
            
            # Use LibreOffice for high-quality PDF conversion
            # This preserves formatting, fonts, images, and layouts perfectly
            try:
                result = subprocess.run(
                    [
                        'libreoffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', temp_dir,
                        temp_input
                    ],
                    capture_output=True,
                    timeout=LIBREOFFICE_TIMEOUT,
                    text=True,
                    check=True
                )
                
                # LibreOffice creates output with same name but .pdf extension
                temp_output = os.path.join(temp_dir, f'input.pdf')
                
                if not os.path.exists(temp_output):
                    logger.error(f"LibreOffice conversion failed. Output not found. stderr: {result.stderr}")
                    return jsonify({'error': 'Conversion failed. Output file not generated.'}), 500
                
                output_filename = original_filename.rsplit('.', 1)[0] + '.pdf'
                
                logger.info(f"Successfully converted {original_filename} to PDF")
                
                return send_file(
                    temp_output,
                    mimetype='application/pdf',
                    as_attachment=True,
                    download_name=output_filename
                )
                
            except subprocess.TimeoutExpired:
                logger.error(f"Conversion timeout for {original_filename}")
                return jsonify({'error': 'Conversion timeout. File may be too large or complex.'}), 504
            
            except subprocess.CalledProcessError as e:
                logger.error(f"LibreOffice error: {e.stderr}")
                return jsonify({'error': 'Conversion failed', 'details': e.stderr}), 500
        
        elif target_format == 'txt':
            # Extract text from PowerPoint
            prs = Presentation(temp_input)
            text_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_content.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_content.append(shape.text)
                text_content.append("")  # Empty line between slides
            
            output_text = '\n'.join(text_content)
            output_buffer = io.BytesIO(output_text.encode('utf-8'))
            output_buffer.seek(0)
            
            output_filename = original_filename.rsplit('.', 1)[0] + '.txt'
            
            logger.info(f"Successfully extracted text from {original_filename}")
            
            return send_file(
                output_buffer,
                mimetype='text/plain',
                as_attachment=True,
                download_name=output_filename
            )
        
        else:
            return jsonify({'error': f'Unsupported format: {target_format}'}), 400
        
    except Exception as e:
        logger.exception(f"Unexpected error during conversion: {str(e)}")
        return jsonify({'error': 'Internal server error', 'details': str(e)}), 500
    
    finally:
        # Cleanup temporary files
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                logger.error(f"Failed to cleanup temp directory: {e}")

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5021))
    debug_mode = os.environ.get('FLASK_ENV', 'production') != 'production'
    app.run(host='0.0.0.0', port=port, debug=debug_mode)
