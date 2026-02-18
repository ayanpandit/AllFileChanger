"""Text Extractor – extract text from PDF, DOCX, PPTX, TXT."""

from flask import Blueprint, request, send_file, jsonify
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io, gc

bp = Blueprint('text_extractor', __name__)

@bp.route('/extract', methods=['POST'])
def extract_text():
    try:
        if 'file' not in request.files:
            return jsonify(error='No file provided'), 400

        f = request.files['file']
        ext = f.filename.rsplit('.', 1)[-1].lower() if '.' in f.filename else ''
        raw = f.read()
        parts = []

        if ext == 'pdf':
            reader = PdfReader(io.BytesIO(raw))
            for page in reader.pages:
                parts.append(page.extract_text() or '')
            del reader  # MEMORY MANAGEMENT: free reader
        elif ext in ('docx', 'doc'):
            doc = Document(io.BytesIO(raw))
            for p in doc.paragraphs:
                parts.append(p.text)
            del doc  # MEMORY MANAGEMENT: free document
        elif ext in ('pptx', 'ppt'):
            prs = Presentation(io.BytesIO(raw))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        parts.append(shape.text)
            del prs  # MEMORY MANAGEMENT: free presentation
        elif ext == 'txt':
            parts.append(raw.decode('utf-8'))
        else:
            return jsonify(error='Unsupported file type'), 400

        del raw  # MEMORY MANAGEMENT: free raw bytes
        out = io.BytesIO('\n\n'.join(parts).encode('utf-8'))
        del parts  # MEMORY MANAGEMENT: free parts list
        out.seek(0)
        return send_file(out, mimetype='text/plain',
                         as_attachment=True, download_name='extracted_text.txt')
    except Exception as e:
        return jsonify(error='Failed to extract text', details=str(e)), 500
    finally:
        gc.collect()
