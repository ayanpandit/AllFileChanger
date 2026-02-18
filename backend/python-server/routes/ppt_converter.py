"""PPT/PPTX Converter – convert to PDF (via LibreOffice) or TXT."""

from flask import Blueprint, request, send_file, jsonify
from pptx import Presentation
import io, os, tempfile, subprocess, shutil, logging, gc

bp = Blueprint('ppt_converter', __name__)
logger = logging.getLogger(__name__)

def _has_libreoffice():
    try:
        subprocess.run(['libreoffice', '--version'], capture_output=True, timeout=5)
        return True
    except Exception:
        return False

_LO = _has_libreoffice()

@bp.route('/ppt-convert', methods=['POST'])
def convert_ppt():
    temp_dir = None
    try:
        if 'file' not in request.files:
            return jsonify(error='No file provided'), 400

        f = request.files['file']
        fmt = request.form.get('format', 'pdf').lower()
        ext = f.filename.rsplit('.', 1)[-1].lower() if '.' in f.filename else 'pptx'

        temp_dir = tempfile.mkdtemp()
        in_path = os.path.join(temp_dir, f'input.{ext}')
        f.save(in_path)

        if fmt == 'pdf':
            if not _LO:
                return jsonify(error='LibreOffice not available on server'), 503
            subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, in_path],
                capture_output=True, timeout=120, check=True
            )
            out_path = os.path.join(temp_dir, 'input.pdf')
            if not os.path.exists(out_path):
                return jsonify(error='Conversion failed'), 500
            return send_file(out_path, mimetype='application/pdf',
                             as_attachment=True, download_name=f.filename.rsplit('.', 1)[0] + '.pdf')

        elif fmt == 'txt':
            prs = Presentation(in_path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f'--- Slide {i} ---')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        lines.append(shape.text)
                lines.append('')
            out = io.BytesIO('\n'.join(lines).encode('utf-8'))
            out.seek(0)
            return send_file(out, mimetype='text/plain',
                             as_attachment=True, download_name=f.filename.rsplit('.', 1)[0] + '.txt')

        return jsonify(error=f'Unsupported format: {fmt}'), 400
    except subprocess.TimeoutExpired:
        return jsonify(error='Conversion timed out'), 504
    except Exception as e:
        logger.exception('PPT convert failed')
        return jsonify(error='Conversion failed', details=str(e)), 500
    finally:
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
        gc.collect()
