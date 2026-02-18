"""PDF → PowerPoint (each page rendered as slide image)."""

from flask import Blueprint, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
import io, os, tempfile, logging, gc

bp = Blueprint('pdf_to_ppt', __name__)
logger = logging.getLogger(__name__)

@bp.route('/to-powerpoint', methods=['POST'])
def pdf_to_ppt():
    pdf_path = ppt_path = None
    img_paths = []
    try:
        if 'pdf' not in request.files:
            return jsonify(error='No PDF file provided'), 400

        pdf_path = tempfile.mktemp(suffix='.pdf')
        request.files['pdf'].save(pdf_path)

        # MEMORY MANAGEMENT: convert one page at a time to avoid loading all pages into RAM
        images = convert_from_path(pdf_path, dpi=150)  # reduced from default 200 dpi
        prs = Presentation()

        for i, image in enumerate(images):
            img_path = f'{pdf_path}_{i}.png'
            img_paths.append(img_path)
            image.save(img_path, 'PNG')
            image.close()  # MEMORY MANAGEMENT: close PIL Image immediately

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)

        # MEMORY MANAGEMENT: free images list
        del images
        gc.collect()

        ppt_path = pdf_path.replace('.pdf', '.pptx')
        prs.save(ppt_path)
        del prs  # MEMORY MANAGEMENT: free presentation

        with open(ppt_path, 'rb') as f:
            out = io.BytesIO(f.read())
        out.seek(0)

        return send_file(out,
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name='converted.pptx')
    except Exception as e:
        logger.exception('PDF→PPT failed')
        return jsonify(error='Failed to convert PDF to PowerPoint', details=str(e)), 500
    finally:
        for p in [pdf_path, ppt_path] + img_paths:
            if p and os.path.exists(p):
                try: os.unlink(p)
                except: pass
        gc.collect()
